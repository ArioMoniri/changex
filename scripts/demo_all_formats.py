"""Track AI edits across ALL supported formats (.docx / .xlsx / .csv / .pptx).

For each format this builds a small sample, applies a few semantic ops through the
format's adapter, records the provenance journal, writes the format's tracked/overlay
output, and renders the HTML review surfaces. Outputs land in ``examples/out/formats``.

Run:  python scripts/demo_all_formats.py
"""

from __future__ import annotations

from pathlib import Path

import changex_core as cx
from changex_core.adapters import load_adapter
from changex_core.journal.events import utc_now_iso
from changex_core.ops.vocabulary import (
    CellSet,
    FormulaSet,
    RowDelete,
    RowInsert,
    ShapeEdit,
    TextDelete,
    TextReplace,
    target_node_id,
)

_OUT = Path(__file__).resolve().parents[1] / "examples" / "out" / "formats"
_AGENT = "claude-opus-4-8"


def _track(adapter, ops, *, fmt: str, filename: str, base: Path) -> cx.Journal:
    """Apply ops through an adapter while recording the provenance journal."""
    node_map = getattr(adapter, "node_id_map", lambda: {})()
    header = cx.Header.create(
        baseline_sha256=adapter.baseline_sha256(),
        filename=filename,
        doc_format=fmt,
        node_id_map=node_map,
    )
    changex_path = base.with_suffix(".changex")
    if changex_path.exists():
        changex_path.unlink()
    journal = cx.Journal.open(str(changex_path), header=header)
    for op, why in ops:
        adapter.apply(op)
        journal.append(
            op,
            cx.Target(node_id=target_node_id(op) or "", node_kind=op.kind.split(".")[0]),
            cx.Provenance(
                ts=utc_now_iso(),
                session_id=header.session_id,
                agent=_AGENT,
                vendor="anthropic",
                provenance_source="declared",
                rationale=why,
            ),
        )
    return journal


def _write(adapter, journal, base: Path, *, tracked_suffix: str, document_view: bool) -> list[str]:
    out: list[str] = []
    tracked = base.with_name(base.name + "-tracked" + tracked_suffix)
    tracked.write_bytes(adapter.render_tracked())
    out.append(str(tracked))
    review = base.with_name(base.name + "-review.html")
    review.write_text(cx.render_html(journal.active_events(), title=f"ChangeX — {base.name}"), "utf-8")
    out.append(str(review))
    if document_view:
        doc = base.with_name(base.name + "-document.html")
        doc.write_text(
            cx.render_document_html(str(tracked), title=f"ChangeX — {base.name}", events=journal.active_events()),
            "utf-8",
        )
        out.append(str(doc))
    return out


def demo_docx() -> list[str]:
    from docx import Document

    base = _OUT / "report"
    src = base.with_suffix(".docx")
    doc = Document()
    h = doc.add_paragraph("Quarterly Report")
    h.style = doc.styles["Heading 1"]
    p = doc.add_paragraph("The quick brown fox jumps over the lazy dog.")
    doc.save(str(src))
    adapter = load_adapter(str(src), author=_AGENT)
    body = next(x for x in adapter.to_model().child_paragraphs() if "quick" in x.text())
    ops = [
        (TextReplace(node_id=body.node_id, before="quick", after="swift"), "tighten wording"),
        (TextDelete(node_id=body.node_id, before=" lazy"), "drop filler"),
    ]
    j = _track(adapter, ops, fmt="docx", filename="report.docx", base=base)
    return _write(adapter, j, base, tracked_suffix=".docx", document_view=True)


def demo_xlsx() -> list[str]:
    from openpyxl import Workbook

    base = _OUT / "budget"
    src = base.with_suffix(".xlsx")
    wb = Workbook()
    ws = wb.active
    ws.title = "Q3"
    ws["A1"], ws["B1"], ws["C1"] = "Item", "Units", "Projected"
    ws["A2"], ws["B2"], ws["C2"] = "Widget", 100, "=B2*1.1"
    ws["A3"], ws["B3"] = "Gadget", 200
    wb.save(str(src))
    adapter = load_adapter(str(src), author=_AGENT)
    ops = [
        (CellSet(sheet="Q3", ref="B2", before="100", after="125"), "revise Widget units"),
        (FormulaSet(sheet="Q3", ref="C2", before="=B2*1.1", after="=B2*1.2"), "bump growth assumption"),
        (RowInsert(sheet="Q3", at=3), "insert a new line item"),
    ]
    j = _track(adapter, ops, fmt="xlsx", filename="budget.xlsx", base=base)
    return _write(adapter, j, base, tracked_suffix=".xlsx", document_view=False)


def demo_csv() -> list[str]:
    base = _OUT / "data"
    src = base.with_suffix(".csv")
    src.write_text("name,score\nalice,10\nbob,20\n", encoding="utf-8")
    adapter = load_adapter(str(src), author=_AGENT)
    ops = [
        (CellSet(sheet="data", ref="B2", before="10", after="15"), "correct alice's score"),
        (RowDelete(sheet="data", at=3, value=["bob", "20"]), "remove bob"),
    ]
    j = _track(adapter, ops, fmt="csv", filename="data.csv", base=base)
    return _write(adapter, j, base, tracked_suffix=".html", document_view=False)


def demo_pptx() -> list[str]:
    from pptx import Presentation
    from pptx.util import Inches

    base = _OUT / "deck"
    src = base.with_suffix(".pptx")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
    box.text_frame.text = "The quick brown fox jumps over the lazy dog."
    prs.save(str(src))
    adapter = load_adapter(str(src), author=_AGENT)
    model = adapter.to_model()
    slide0 = model.children[0]
    body = next(s for s in slide0.children if "shape:" in s.node_id and s.children)
    sid = body.node_id.split("shape:")[-1]
    para_id = body.children[0].node_id
    ops = [
        (ShapeEdit(slide=0, shape_id=sid, op={"kind": "text.replace", "node_id": para_id, "before": "quick", "after": "swift"}), "tighten slide copy"),
        (ShapeEdit(slide=0, shape_id=sid, op={"kind": "text.delete", "node_id": para_id, "before": " over the lazy dog"}), "trim"),
    ]
    j = _track(adapter, ops, fmt="pptx", filename="deck.pptx", base=base)
    return _write(adapter, j, base, tracked_suffix=".pptx", document_view=False)


def main() -> None:
    _OUT.mkdir(parents=True, exist_ok=True)
    produced: dict[str, list[str]] = {}
    for name, fn in (("docx", demo_docx), ("xlsx", demo_xlsx), ("csv", demo_csv), ("pptx", demo_pptx)):
        try:
            produced[name] = fn()
            print(f"[{name}] OK")
            for path in produced[name]:
                print(f"    {path}")
        except Exception as exc:  # surface per-format so one failure doesn't hide the rest
            print(f"[{name}] FAILED: {type(exc).__name__}: {exc}")
    print("\nMANIFEST")
    for paths in produced.values():
        for path in paths:
            print(path)


if __name__ == "__main__":
    main()
