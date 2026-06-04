"""Render review surfaces for README screenshots into docs/assets/_render/*.html.

Builds a *longer, realistic* sample per format, applies several tracked edits spread
through it, and writes the HTML review the app shows by default:
- docx -> the in-document outline (render_document_html)
- md / csv -> the inline redline (adapter.render_tracked)
- xlsx / pptx -> the provenance review (render_html over the journal)

A separate step (make_screenshots.sh / headless Chrome) turns each into a PNG.
"""

from __future__ import annotations

from pathlib import Path

import changex_core as cx
from changex_core.adapters import load_adapter
from changex_core.journal.events import utc_now_iso
from changex_core.ops.vocabulary import (
    CellSet,
    FormulaSet,
    RowInsert,
    ShapeEdit,
    TextDelete,
    TextInsert,
    TextReplace,
    target_node_id,
)

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "assets" / "_render"
AGENT = "claude-opus-4-8"


def _track(adapter, ops, *, fmt, filename, base):
    nim = getattr(adapter, "node_id_map", lambda: {})()
    header = cx.Header.create(
        baseline_sha256=adapter.baseline_sha256(), filename=filename, doc_format=fmt, node_id_map=nim
    )
    cxp = base.with_suffix(".changex")
    if cxp.exists():
        cxp.unlink()
    j = cx.Journal.open(str(cxp), header=header)
    for op, why in ops:
        adapter.apply(op)
        j.append(
            op,
            cx.Target(node_id=target_node_id(op) or "", node_kind=op.kind.split(".")[0]),
            cx.Provenance(
                ts=utc_now_iso(), session_id=header.session_id, agent=AGENT,
                vendor="anthropic", provenance_source="declared", rationale=why,
            ),
        )
    return j


def docx() -> None:
    from docx import Document

    base = OUT / "report"
    src = base.with_suffix(".docx")
    d = Document()
    d.add_paragraph("Q3 Platform Status Report").style = d.styles["Title"]
    d.add_paragraph(
        "Prepared for the executive review on October 14. This memo summarizes "
        "progress, the key results, and the open risks heading into Q4."
    )
    d.add_paragraph("1. Executive summary").style = d.styles["Heading 1"]
    d.add_paragraph(
        "The platform shipped three major features this quarter and grew weekly "
        "active users by a healthy margin. Reliability held steady despite the quick "
        "increase in traffic, and the team closed the most pressing security findings "
        "ahead of schedule."
    )
    d.add_paragraph("2. Key results").style = d.styles["Heading 1"]
    d.add_paragraph(
        "Adoption exceeded the lazy targets we set in July, with the redesigned "
        "onboarding flow driving most of the lift. Support volume fell as the "
        "documentation matured."
    )
    d.add_paragraph(
        "Infrastructure cost was roughly flat, and the migration to the new datastore "
        "finished without a single customer-visible incident."
    )
    d.add_paragraph("3. Risks and mitigations").style = d.styles["Heading 1"]
    d.add_paragraph(
        "The biggest risk remains our dependency on a single upstream vendor for "
        "document conversion. A backup path is being scoped."
    )
    d.add_paragraph("Prepared by the analytics team.")
    d.save(str(src))

    a = load_adapter(str(src), author=AGENT)
    paras = a.to_model().child_paragraphs()

    def find(sub):
        return next(p for p in paras if sub in p.text())

    ops = [
        (TextReplace(node_id=find("quick increase").node_id, before="quick", after="rapid"), "tighten wording"),
        (TextDelete(node_id=find("three major features").node_id, before=" this quarter"), "trim redundancy"),
        (TextReplace(node_id=find("lazy targets").node_id, before="lazy", after="conservative"), "more precise"),
        (TextReplace(node_id=find("Infrastructure cost").node_id, before="roughly flat", after="down 4%"), "use the real number"),
        (TextInsert(node_id=find("being scoped").node_id, before_anchor="being scoped", text=" and lands in early Q4"), "add the commitment"),
    ]
    j = _track(a, ops, fmt="docx", filename="report.docx", base=base)
    tracked = base.with_name("report.tracked.docx")
    tracked.write_bytes(a.render_tracked())
    (OUT / "docx.html").write_text(
        cx.render_document_html(str(tracked), title="Q3 Platform Status Report — ChangeX", events=j.active_events()),
        "utf-8",
    )


def md() -> None:
    base = OUT / "guide"
    src = base.with_suffix(".md")
    src.write_text(
        "# Onboarding guide\n\n"
        "Welcome aboard! This guide walks you through your first day.\n\n"
        "## 1. Set up your laptop\n\n"
        "Install the quick tools from the internal portal and request access to the repos.\n\n"
        "## 2. Join the standups\n\n"
        "Standup is at 10am every day. Bring your blockers and a coffee.\n\n"
        "## 3. Ship something small\n\n"
        "Pick a good-first-issue and open a pull request by the end of week one.\n",
        encoding="utf-8",
    )
    a = load_adapter(str(src), author=AGENT)
    blocks = a.to_model().children

    def block(sub):
        return next(b for b in blocks if sub in b.text())

    ops = [
        (TextReplace(node_id=block("quick tools").node_id, before="quick", after="required"), "be specific"),
        (TextReplace(node_id=block("10am every day").node_id, before="10am", after="9:30am"), "fix the time"),
        (TextInsert(node_id=block("good-first-issue").node_id, before_anchor="week one", text=" — your buddy will review it"), "add support"),
    ]
    _track(a, ops, fmt="md", filename="guide.md", base=base)
    (OUT / "md.html").write_text(a.render_tracked().decode("utf-8"), "utf-8")


def csv() -> None:
    base = OUT / "headcount"
    src = base.with_suffix(".csv")
    src.write_text(
        "team,headcount,open_roles\n"
        "Platform,12,2\n"
        "Growth,8,1\n"
        "Design,5,0\n"
        "Data,6,3\n",
        encoding="utf-8",
    )
    a = load_adapter(str(src), author=AGENT)
    ops = [
        (CellSet(sheet="headcount", ref="B2", before="12", after="14"), "Platform hired two"),
        (CellSet(sheet="headcount", ref="C5", before="3", after="1"), "Data filled two roles"),
        (RowInsert(sheet="headcount", at=4), "add the new Research team"),
    ]
    _track(a, ops, fmt="csv", filename="headcount.csv", base=base)
    (OUT / "csv.html").write_text(a.render_tracked().decode("utf-8"), "utf-8")


def xlsx() -> None:
    from openpyxl import Workbook

    base = OUT / "budget"
    src = base.with_suffix(".xlsx")
    wb = Workbook()
    ws = wb.active
    ws.title = "Q4"
    ws["A1"], ws["B1"], ws["C1"] = "Line item", "Q3 actual", "Q4 plan"
    rows = [("Salaries", 420, "=B2*1.05"), ("Cloud", 95, "=B3*1.1"), ("Travel", 18, "=B4*1.0"), ("Tooling", 30, "=B5*1.2")]
    for i, (name, q3, plan) in enumerate(rows, start=2):
        ws[f"A{i}"], ws[f"B{i}"], ws[f"C{i}"] = name, q3, plan
    wb.save(str(src))
    a = load_adapter(str(src), author=AGENT)
    ops = [
        (CellSet(sheet="Q4", ref="B3", before="95", after="120"), "cloud spend rose"),
        (FormulaSet(sheet="Q4", ref="C3", before="=B3*1.1", after="=B3*1.15"), "higher growth assumption"),
        (RowInsert(sheet="Q4", at=6), "add a Security line item"),
    ]
    j = _track(a, ops, fmt="xlsx", filename="budget.xlsx", base=base)
    (OUT / "xlsx.html").write_text(
        cx.render_html(j.active_events(), title="Q4 Budget — ChangeX review"), "utf-8"
    )


def pptx() -> None:
    from pptx import Presentation
    from pptx.util import Inches

    base = OUT / "deck"
    src = base.with_suffix(".pptx")
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(8), Inches(1)).text_frame.text = "Q3 in review"
    s.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8), Inches(3)).text_frame.text = (
        "We shipped quickly and grew the lazy way — by word of mouth."
    )
    prs.save(str(src))
    a = load_adapter(str(src), author=AGENT)
    slide0 = a.to_model().children[0]
    body = next(sh for sh in slide0.children if "shape:" in sh.node_id and sh.children and "lazy" in sh.children[0].text())
    sid = body.node_id.split("shape:")[-1]
    pid = body.children[0].node_id
    ops = [
        (ShapeEdit(slide=0, shape_id=sid, op={"kind": "text.replace", "node_id": pid, "before": "quickly", "after": "deliberately"}), "tone"),
        (ShapeEdit(slide=0, shape_id=sid, op={"kind": "text.delete", "node_id": pid, "before": " the lazy way"}), "trim"),
    ]
    j = _track(a, ops, fmt="pptx", filename="deck.pptx", base=base)
    (OUT / "pptx.html").write_text(
        cx.render_html(j.active_events(), title="Q3 deck — ChangeX review"), "utf-8"
    )


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    for name, fn in (("docx", docx), ("md", md), ("csv", csv), ("xlsx", xlsx), ("pptx", pptx)):
        try:
            fn()
            print(f"[{name}] OK -> {OUT / (name + '.html')}")
        except Exception as exc:
            print(f"[{name}] FAILED: {type(exc).__name__}: {exc}")


if __name__ == "__main__":
    main()
