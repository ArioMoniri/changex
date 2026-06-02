"""PptxAdapter + pptx review-overlay tests against a real generated ``.pptx``.

These cover the M4 pptx scope via the public ``changex_core`` API:

* stable, positional node_ids for slides / shapes / paragraphs, with out-of-scope
  shapes (a table) surfaced as **opaque, untracked** nodes;
* the in-scope ops apply: ``shape.edit`` text replace/insert/delete/style on a
  text frame, plus ``slide.insert`` / ``slide.delete`` (add / remove / reorder);
* boundary guards — before-mismatch, oversized op, missing node, and the honest
  **refusal** to edit an opaque (chart/table/picture) shape;
* the journal of those ops verifies and replays back to the live model, including
  a MIDDLE-op revert;
* ``render_tracked`` yields a re-openable deck reflecting the edits; and
* the review overlay (:mod:`changex_core.render.pptx`) is **non-destructive** —
  it never mutates the shipped slides and appends a "Revisions" summary slide.
"""

from __future__ import annotations

import io
from pathlib import Path

import pytest

import changex_core as cx
from changex_core.journal.events import utc_now_iso

pytest.importorskip("pptx", reason="python-pptx is required for the pptx adapter")

from pptx import Presentation  # noqa: E402  (after importorskip)
from pptx.util import Inches  # noqa: E402

from changex_core.adapters import adapter_class_for, load_adapter  # noqa: E402
from changex_core.adapters.pptx_adapter import (  # noqa: E402
    OpaqueNodeError,
    PptxAdapter,
)
from changex_core.render.pptx import render_review_overlay, save_review_overlay  # noqa: E402


# --------------------------------------------------------------------------- #
# fixtures
# --------------------------------------------------------------------------- #
@pytest.fixture()
def sample_pptx(tmp_path: Path) -> Path:
    """A freshly generated sample ``.pptx``: 2 slides; a textbox; a table (opaque)."""
    prs = Presentation()
    s0 = prs.slides.add_slide(prs.slide_layouts[0])  # title + subtitle
    s0.shapes.title.text = "Quarterly Review"
    s0.placeholders[1].text = "The quick brown fox jumps over the lazy dog"
    s1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    tb = s1.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    tb.text_frame.text = "draft notes line"
    # A table is OUT OF SCOPE — it must show up as an opaque, untracked shape.
    s1.shapes.add_table(2, 2, Inches(1), Inches(3), Inches(4), Inches(1))
    out = tmp_path / "sample.pptx"
    prs.save(str(out))
    return out


def _prov(session_id: str) -> cx.Provenance:
    return cx.Provenance(
        ts=utc_now_iso(),
        session_id=session_id,
        agent="claude-opus-4-8",
        vendor="anthropic",
        provenance_source="declared",
    )


def _slides(adapter: PptxAdapter) -> list[cx.Node]:
    return list(adapter.to_model().children)


def _body_shape_id(adapter: PptxAdapter, slide_idx: int, text_contains: str) -> str:
    """Return the shape_id (last path segment) of the text shape containing text."""
    slide = _slides(adapter)[slide_idx]
    for shape in slide.children:
        if shape.attrs.get("opaque"):
            continue
        if any(text_contains in str(p.value) for p in shape.children):
            return shape.node_id.split("shape:")[-1]
    raise AssertionError(f"no shape on slide {slide_idx} containing {text_contains!r}")


# --------------------------------------------------------------------------- #
# Addressing: positional node_ids + opaque out-of-scope shapes
# --------------------------------------------------------------------------- #
def test_model_has_positional_slide_shape_paragraph_ids(sample_pptx: Path) -> None:
    adapter = PptxAdapter.load(str(sample_pptx))
    model = adapter.to_model()
    slides = list(model.children)
    assert [s.node_id for s in slides] == ["slide[0]", "slide[1]"]
    assert slides[0].attrs["title"] == "Quarterly Review"
    # Every in-scope shape node_id is slide[i]/shape:<id>; paragraphs add /p[j].
    body = next(
        sh for sh in slides[0].children
        if not sh.attrs["opaque"] and any("fox" in str(p.value) for p in sh.children)
    )
    assert body.node_id.startswith("slide[0]/shape:")
    assert body.children[0].node_id == f"{body.node_id}/p[0]"
    assert "fox" in str(body.children[0].value)


def test_table_shape_is_opaque_and_untracked(sample_pptx: Path) -> None:
    adapter = PptxAdapter.load(str(sample_pptx))
    slide1 = _slides(adapter)[1]
    opaque = [sh for sh in slide1.children if sh.attrs.get("opaque")]
    assert opaque, "the table should surface as an opaque shape"
    assert opaque[0].attrs["shape_kind"] == "table"
    assert opaque[0].children == []  # no tracked paragraphs


def test_node_id_map_records_native_slide_carriers(sample_pptx: Path) -> None:
    adapter = PptxAdapter.load(str(sample_pptx))
    carriers = adapter.node_id_map()
    assert set(carriers) == {"slide[0]", "slide[1]"}
    # carrier is the native pptx slide_id (a stable integer string)
    assert all(v.isdigit() for v in carriers.values())


# --------------------------------------------------------------------------- #
# shape.edit text ops + slide.insert / slide.delete
# --------------------------------------------------------------------------- #
def test_shape_edit_text_replace_insert_delete(sample_pptx: Path) -> None:
    adapter = PptxAdapter.load(str(sample_pptx))
    sid = _body_shape_id(adapter, 0, "fox")
    para_id = f"slide[0]/shape:{sid}/p[0]"

    adapter.apply(cx.ShapeEdit(slide=0, shape_id=sid, op={
        "kind": "text.replace", "node_id": para_id, "before": "quick", "after": "swift",
    }))
    assert "swift" in str(adapter.resolve(para_id).value)
    assert "quick" not in str(adapter.resolve(para_id).value)

    adapter.apply(cx.ShapeEdit(slide=0, shape_id=sid, op={
        "kind": "text.insert", "node_id": para_id, "before_anchor": "fox", "text": " (red)",
    }))
    assert "fox (red)" in str(adapter.resolve(para_id).value)

    adapter.apply(cx.ShapeEdit(slide=0, shape_id=sid, op={
        "kind": "text.delete", "node_id": para_id, "before": " over the lazy dog",
    }))
    assert "lazy dog" not in str(adapter.resolve(para_id).value)


def test_shape_edit_style_change_validates_before(sample_pptx: Path) -> None:
    adapter = PptxAdapter.load(str(sample_pptx))
    sid = _body_shape_id(adapter, 0, "fox")
    para_id = f"slide[0]/shape:{sid}/p[0]"
    adapter.apply(cx.ShapeEdit(slide=0, shape_id=sid, op={
        "kind": "style.change", "node_id": para_id, "style": "Heading", "before": "",
    }))
    assert adapter.resolve(para_id).attrs.get("style") == "Heading"
    # A wrong `before` style is refused.
    with pytest.raises(cx.BeforeMismatchError):
        adapter.apply(cx.ShapeEdit(slide=0, shape_id=sid, op={
            "kind": "style.change", "node_id": para_id, "style": "Body", "before": "Normal",
        }))


def test_slide_insert_and_delete_reorder(sample_pptx: Path) -> None:
    adapter = PptxAdapter.load(str(sample_pptx))
    assert len(_slides(adapter)) == 2
    adapter.apply(cx.SlideInsert(at=1, value={"layout": "Blank", "title": "Inserted"}))
    slides = _slides(adapter)
    assert len(slides) == 3
    assert slides[1].attrs["title"] == "Inserted"
    assert slides[1].attrs["inserted"] is True
    # delete the original slide-1 (now at index 2)
    adapter.apply(cx.SlideDelete(at=2, value={"title": ""}))
    titles = [s.attrs["title"] for s in _slides(adapter)]
    assert titles == ["Quarterly Review", "Inserted"]


# --------------------------------------------------------------------------- #
# Boundary guards: opaque refusal, before-mismatch, oversized, missing node
# --------------------------------------------------------------------------- #
def test_editing_opaque_shape_is_refused(sample_pptx: Path) -> None:
    adapter = PptxAdapter.load(str(sample_pptx))
    table = next(sh for sh in _slides(adapter)[1].children if sh.attrs.get("opaque"))
    table_id = table.node_id.split("shape:")[-1]
    with pytest.raises(OpaqueNodeError):
        adapter.apply(cx.ShapeEdit(slide=1, shape_id=table_id, op={
            "kind": "text.replace", "node_id": "", "before": "a", "after": "b",
        }))


def test_before_mismatch_and_oversized_and_missing(sample_pptx: Path) -> None:
    adapter = PptxAdapter.load(str(sample_pptx))
    sid = _body_shape_id(adapter, 0, "fox")
    para_id = f"slide[0]/shape:{sid}/p[0]"
    # before not present -> BeforeMismatchError
    with pytest.raises(cx.BeforeMismatchError):
        adapter.apply(cx.ShapeEdit(slide=0, shape_id=sid, op={
            "kind": "text.replace", "node_id": para_id, "before": "absent", "after": "x",
        }))
    # rewriting the whole paragraph -> OversizedOpError (split_required)
    whole = str(adapter.resolve(para_id).value)
    with pytest.raises(cx.OversizedOpError) as exc:
        adapter.apply(cx.ShapeEdit(slide=0, shape_id=sid, op={
            "kind": "text.replace", "node_id": para_id, "before": whole, "after": "z",
        }))
    assert "split_required" in str(exc.value)
    # missing slide / shape -> NodeNotFoundError
    with pytest.raises(cx.NodeNotFoundError):
        adapter.apply(cx.SlideDelete(at=99, value={}))
    with pytest.raises(cx.NodeNotFoundError):
        adapter.apply(cx.ShapeEdit(slide=0, shape_id="nope", op={
            "kind": "text.delete", "node_id": "", "before": "x",
        }))


# --------------------------------------------------------------------------- #
# Registry + render round-trip
# --------------------------------------------------------------------------- #
def test_registry_resolves_pptx_to_adapter(sample_pptx: Path) -> None:
    assert adapter_class_for(str(sample_pptx)) is PptxAdapter
    adapter = load_adapter(str(sample_pptx), author="claude-opus-4-8")
    assert isinstance(adapter, PptxAdapter)


def test_render_tracked_reopens_with_edits_applied(sample_pptx: Path, tmp_path: Path) -> None:
    adapter = PptxAdapter.load(str(sample_pptx), author="claude-opus-4-8")
    sid = _body_shape_id(adapter, 0, "fox")
    adapter.apply(cx.ShapeEdit(slide=0, shape_id=sid, op={
        "kind": "text.replace", "node_id": f"slide[0]/shape:{sid}/p[0]",
        "before": "quick", "after": "swift",
    }))
    adapter.apply(cx.SlideInsert(at=1, value={"layout": "Blank", "title": "Added"}))
    adapter.apply(cx.SlideDelete(at=2, value={"title": ""}))
    out = tmp_path / "tracked.pptx"
    adapter.save(str(out))

    prs = Presentation(str(out))
    slides = list(prs.slides)
    assert len(slides) == 2  # original-1 deleted, 1 inserted
    texts = [
        sh.text_frame.text for sl in slides for sh in sl.shapes if sh.has_text_frame
    ]
    assert any("swift" in t for t in texts)
    assert all("quick" not in t for t in texts)
    assert any("Added" in t for t in texts)


# --------------------------------------------------------------------------- #
# Journal integration: verify + replay + middle revert
# --------------------------------------------------------------------------- #
def test_journal_over_pptx_verifies_replays_and_reverts(
    sample_pptx: Path, tmp_path: Path
) -> None:
    adapter = PptxAdapter.load(str(sample_pptx), author="claude-opus-4-8")
    baseline = cx.Node.from_dict(adapter.to_model().to_dict())
    sid = _body_shape_id(adapter, 0, "fox")
    para_id = f"slide[0]/shape:{sid}/p[0]"

    header = cx.Header.create(
        baseline_sha256=adapter.baseline_sha256(),
        filename=sample_pptx.name,
        doc_format="pptx",
        node_id_map=adapter.node_id_map(),
    )
    journal = cx.Journal.open(str(tmp_path / "session.changex"), header=header)
    sess = journal.header.session_id

    def record(op: cx.Op) -> cx.Event:
        adapter.apply(op)
        nid = cx.target_node_id(op) or ""
        node = adapter.resolve(nid)
        target = cx.Target(
            node_id=nid,
            node_kind=(node.node_kind.value if node else "document"),
            path=(node.path if node else ""),
        )
        return journal.append(op, target, _prov(sess))

    record(cx.ShapeEdit(slide=0, shape_id=sid, op={
        "kind": "text.replace", "node_id": para_id, "before": "quick", "after": "swift",
    }))
    e2 = record(cx.ShapeEdit(slide=0, shape_id=sid, op={
        "kind": "text.replace", "node_id": para_id, "before": "lazy", "after": "sleepy",
    }))
    record(cx.SlideInsert(at=1, value={"layout": "Blank", "title": "Summary"}))

    # Chain verifies and replay reproduces the live model.
    assert journal.verify().ok
    fresh = load_adapter(str(sample_pptx), author="claude-opus-4-8")
    replayed = journal.replay(fresh, baseline)
    live = {s.node_id: s.attrs.get("title") for s in adapter.to_model().children}
    assert {s.node_id: s.attrs.get("title") for s in replayed.children} == live

    # Reject the MIDDLE op (lazy -> sleepy): 'swift' kept, 'lazy' restored.
    journal.revert(e2.op_id)
    fresh2 = load_adapter(str(sample_pptx), author="claude-opus-4-8")
    after = journal.replay(fresh2, baseline)
    body_text = str(after.find(para_id).value)
    assert "swift" in body_text  # first edit survives
    assert "lazy" in body_text and "sleepy" not in body_text  # middle reverted


# --------------------------------------------------------------------------- #
# Review overlay: NON-DESTRUCTIVE separate copy + generated Revisions slide
# --------------------------------------------------------------------------- #
def test_review_overlay_is_non_destructive_and_adds_summary(
    sample_pptx: Path, tmp_path: Path
) -> None:
    adapter = PptxAdapter.load(str(sample_pptx), author="claude-opus-4-8")
    sid = _body_shape_id(adapter, 0, "fox")
    header = cx.Header.create(
        baseline_sha256=adapter.baseline_sha256(), filename=sample_pptx.name,
        doc_format="pptx",
    )
    journal = cx.Journal.open(str(tmp_path / "rev.changex"), header=header)
    op = cx.ShapeEdit(slide=0, shape_id=sid, op={
        "kind": "text.replace", "node_id": f"slide[0]/shape:{sid}/p[0]",
        "before": "quick", "after": "swift",
    })
    adapter.apply(op)
    journal.append(op, cx.Target(node_id="slide[0]", node_kind="document"), _prov(journal.header.session_id))

    base = Presentation(str(sample_pptx))
    overlay_bytes = render_review_overlay(str(sample_pptx), journal.active_events())
    overlay = Presentation(io.BytesIO(overlay_bytes))

    base_slides = list(base.slides)
    ov_slides = list(overlay.slides)
    # overlay = baseline + exactly one generated "Revisions" summary slide
    assert len(ov_slides) == len(base_slides) + 1

    # Original slides' TEXT is unchanged (non-destructive): the overlay must NOT
    # carry the 'swift' edit — it annotates, it does not apply.
    orig_text = " ".join(
        sh.text_frame.text for sl in base_slides for sh in sl.shapes if sh.has_text_frame
    )
    ov_orig_text = " ".join(
        sh.text_frame.text
        for sl in ov_slides[: len(base_slides)]
        for sh in sl.shapes
        if sh.has_text_frame
    )
    assert "quick" in orig_text and "swift" not in orig_text
    assert ov_orig_text == orig_text  # byte-for-text identical original slides

    # The appended summary slide names the revision + the honesty disclaimer.
    summary = ov_slides[-1]
    summary_text = " ".join(
        sh.text_frame.text for sh in summary.shapes if sh.has_text_frame
    )
    assert "Revisions" in summary_text
    assert "no native track-changes" in summary_text.lower()
    assert "swift" in summary_text  # the change is described in the summary


def test_save_review_overlay_writes_pptx(sample_pptx: Path, tmp_path: Path) -> None:
    header = cx.Header.create(
        baseline_sha256="0" * 64, filename=sample_pptx.name, doc_format="pptx"
    )
    journal = cx.Journal.open(str(tmp_path / "ov.changex"), header=header)
    out = tmp_path / "review.pptx"
    written = save_review_overlay(str(sample_pptx), journal.active_events(), str(out))
    assert Path(written).exists()
    # Re-opens cleanly with the summary slide (no events -> "No tracked changes").
    prs = Presentation(written)
    last = list(prs.slides)[-1]
    text = " ".join(sh.text_frame.text for sh in last.shapes if sh.has_text_frame)
    assert "No tracked changes" in text
