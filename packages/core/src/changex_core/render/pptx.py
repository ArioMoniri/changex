"""Non-destructive pptx review overlay: a SEPARATE annotated review copy.

There is **no native track-changes in PowerPoint.** This module therefore does
*not* and *cannot* produce a deck a reviewer "accepts/rejects" in PowerPoint —
that claim would be dishonest (see ``docs/FIDELITY.md``). Instead it builds a
**separate review copy** of the deck that:

* leaves every shipped slide's content **unmutated** — the original slides are
  copied through untouched; nothing is overwritten in place; and
* **appends** a generated "Revisions" summary slide that lists every active
  journal op (kind, target, change, author), plus optionally a per-slide speaker
  **note** flagging which slides carry tracked changes.

The authoritative change record is the ``.changex`` journal and the HTML review
surface (:mod:`changex_core.render.html`). This overlay is a convenience so a
reader can see *that* and *what* changed without leaving PowerPoint — it never
asserts a resolution the host application did not perform.

The module degrades gracefully: it imports ``pptx`` lazily so importing
``changex_core`` never requires python-pptx, and the public functions raise a
clear error only when actually invoked without it installed.
"""

from __future__ import annotations

import io
from typing import Iterable, Optional

from changex_core.journal.events import Event
from changex_core.paths import safe_path

_TITLE = "Revisions (ChangeX)"
_DISCLAIMER = (
    "PowerPoint has no native track-changes. This is a non-destructive review "
    "copy: original slides are unchanged; the .changex journal is authoritative."
)


def _change_summary(event: Event) -> str:
    """Return a one-line human description of an op for the summary slide."""
    op = event.op
    kind = str(op.get("kind", "?"))
    if kind == "slide.insert":
        val = op.get("value", {}) or {}
        title = val.get("title") or val.get("layout") or ""
        return f"insert slide at {op.get('at')}" + (f": {title}" if title else "")
    if kind == "slide.delete":
        val = op.get("value", {}) or {}
        title = val.get("title") or ""
        return f"delete slide at {op.get('at')}" + (f": {title}" if title else "")
    if kind == "shape.edit":
        inner = op.get("op", {}) or {}
        ikind = str(inner.get("kind", "?"))
        if ikind == "text.replace":
            return f"text {inner.get('before')!r} -> {inner.get('after')!r}"
        if ikind == "text.insert":
            anchor = inner.get("before_anchor")
            where = f" after {anchor!r}" if anchor else " (append)"
            return f"insert {inner.get('text')!r}{where}"
        if ikind == "text.delete":
            return f"delete {inner.get('before')!r}"
        if ikind == "style.change":
            return f"style {inner.get('before')} -> {inner.get('style')}"
        return f"shape edit ({ikind})"
    return kind


def _provenance_label(event: Event) -> str:
    prov = event.provenance
    who = prov.agent or "unknown"
    return f"{who} ({prov.provenance_source})"


def _summary_lines(events: list[Event]) -> list[str]:
    """Return the body lines for the Revisions summary slide."""
    if not events:
        return ["No tracked changes."]
    lines: list[str] = []
    for event in events:
        target = event.target.node_id or "?"
        lines.append(
            f"seq {event.seq} - {event.op.get('kind', '?')} @ {target} - "
            f"{_change_summary(event)} [{_provenance_label(event)}]"
        )
    return lines


def render_review_overlay(
    baseline_pptx: str,
    events: Iterable[Event],
    *,
    title: str = _TITLE,
    annotate_notes: bool = True,
) -> bytes:
    """Build a non-destructive review copy of ``baseline_pptx`` and return bytes.

    The original deck is opened and copied through with **no slide content
    mutated**. A "Revisions" summary slide listing every event is appended; when
    ``annotate_notes`` is set, slides referenced by an op also get a speaker note
    flagging that they carry tracked changes. The function never claims a native
    accept/reject — it only surfaces the journal inside the deck.

    Args:
        baseline_pptx: Path to the original (unmodified) ``.pptx``.
        events: The journal events (typically ``journal.active_events()``).
        title: Heading for the generated summary slide.
        annotate_notes: If ``True``, add a per-slide speaker note on changed slides.

    Returns:
        The review-copy ``.pptx`` bytes.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt

    resolved = safe_path(baseline_pptx, must_exist=True, allow_suffixes=(".pptx",))
    prs = Presentation(str(resolved))
    event_list = list(events)

    if annotate_notes:
        _annotate_slide_notes(prs, event_list)

    _append_summary_slide(prs, event_list, title=title, Inches=Inches, Pt=Pt)

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def _changed_slide_indices(events: list[Event]) -> dict[int, list[Event]]:
    """Map a 0-based slide index to the events that touch it."""
    by_slide: dict[int, list[Event]] = {}
    for event in events:
        op = event.op
        kind = str(op.get("kind", ""))
        idx: Optional[int] = None
        if kind in ("slide.insert", "slide.delete"):
            at = op.get("at")
            idx = int(at) if isinstance(at, int) else None
        elif kind == "shape.edit":
            sl = op.get("slide")
            idx = int(sl) if isinstance(sl, int) else None
        if idx is not None:
            by_slide.setdefault(idx, []).append(event)
    return by_slide


def _annotate_slide_notes(prs: object, events: list[Event]) -> None:
    """Add a speaker note to each slide that has tracked changes (non-destructive)."""
    slides = list(prs.slides)  # type: ignore[attr-defined]
    for idx, slide_events in _changed_slide_indices(events).items():
        if idx < 0 or idx >= len(slides):
            continue
        slide = slides[idx]
        notes = slide.notes_slide.notes_text_frame
        change_lines = "\n".join(
            f"- {_change_summary(e)} [{_provenance_label(e)}]" for e in slide_events
        )
        marker = f"[ChangeX] {len(slide_events)} tracked change(s):\n{change_lines}"
        existing = notes.text
        notes.text = (existing + "\n" + marker) if existing else marker


def _append_summary_slide(
    prs: object, events: list[Event], *, title: str, Inches: object, Pt: object
) -> None:
    """Append the generated "Revisions" summary slide listing all events."""
    layouts = prs.slide_layouts  # type: ignore[attr-defined]
    layout = layouts[6] if len(layouts) > 6 else layouts[-1]
    slide = prs.slides.add_slide(layout)  # type: ignore[attr-defined]

    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)  # type: ignore[operator]
    )
    title_tf = title_box.text_frame
    title_tf.text = title
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.size = Pt(24)  # type: ignore[operator]

    disclaimer_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.1), Inches(9), Inches(0.6)  # type: ignore[operator]
    )
    disclaimer_tf = disclaimer_box.text_frame
    disclaimer_tf.word_wrap = True
    disclaimer_tf.text = _DISCLAIMER
    disclaimer_tf.paragraphs[0].font.italic = True
    disclaimer_tf.paragraphs[0].font.size = Pt(10)  # type: ignore[operator]

    body_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.9), Inches(9), Inches(5)  # type: ignore[operator]
    )
    body_tf = body_box.text_frame
    body_tf.word_wrap = True
    lines = _summary_lines(events)
    body_tf.text = lines[0]
    body_tf.paragraphs[0].font.size = Pt(12)  # type: ignore[operator]
    for line in lines[1:]:
        para = body_tf.add_paragraph()
        para.text = line
        para.font.size = Pt(12)  # type: ignore[operator]


def save_review_overlay(
    baseline_pptx: str,
    events: Iterable[Event],
    out_path: str,
    *,
    title: str = _TITLE,
    annotate_notes: bool = True,
) -> str:
    """Render :func:`render_review_overlay` and write it to ``out_path``.

    ``out_path`` is sanitized and must carry a ``.pptx`` suffix. Returns the
    resolved output path as a string.
    """
    resolved = safe_path(out_path, allow_suffixes=(".pptx",))
    resolved.parent.mkdir(parents=True, exist_ok=True)
    data = render_review_overlay(
        baseline_pptx, events, title=title, annotate_notes=annotate_notes
    )
    resolved.write_bytes(data)
    return str(resolved)


__all__ = ["render_review_overlay", "save_review_overlay"]
