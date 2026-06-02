"""The pptx adapter: load -> model -> apply v0.2 slide/text ops -> render bytes.

Scope (deliberately narrow, per review): this adapter tracks **slide add /
delete / reorder** (``slide.insert`` / ``slide.delete``) and **text-frame edits**
delivered as a ``shape.edit`` whose nested ``op`` is a docx text op
(``text.insert`` / ``text.delete`` / ``text.replace`` / ``style.change``)
addressed at a paragraph inside a text frame.

Everything else a slide can hold — charts, tables, SmartArt, pictures, grouped
shapes — is **out of scope** and represented as **opaque, untracked** nodes
(``attrs['opaque']=True``). Ops targeting them raise :class:`OpaqueNodeError`
rather than silently mangling content we cannot redline (a ``docs/FIDELITY.md``
honesty rule made mechanical).

There is **no native track-changes in pptx**, so ``render_tracked`` returns the
*edited deck* bytes; the non-destructive review overlay (a generated "Revisions"
summary slide) lives in :mod:`changex_core.render.pptx`.

Identity: internal state is keyed on each slide's native ``slide_id`` (stable
across reorder/insert/delete). The positional node_ids the model exposes
(``slide[i]`` / ``slide[i]/shape:<id>`` / ``slide[i]/shape:<id>/p[j]``) are
recomputed from current order, matching the frozen positional addresses.
"""

from __future__ import annotations

import io
import re
from dataclasses import dataclass, field
from typing import Any, Optional

from changex_core.adapters.base import (
    AdapterError,
    BeforeMismatchError,
    DocumentAdapter,
    NodeNotFoundError,
    OversizedOpError,
)
from changex_core.journal.canonical import sha256_hex
from changex_core.model.nodes import Node, NodeKind
from changex_core.ops.vocabulary import Op, ShapeEdit, SlideDelete, SlideInsert
from changex_core.paths import safe_path

MAX_OP_FRACTION = 0.5  # an op may not rewrite > this fraction of a paragraph
DEFAULT_AUTHOR = "ChangeX agent"
DEFAULT_DATE = "2026-06-02T00:00:00Z"
_BLANK_LAYOUT_INDEX = 6  # blank layout for minted slides (template-agnostic)
_PARA_NODE_RE = re.compile(r"^slide\[(\d+)\]/shape:([^/]+)/p\[(\d+)\]$")  # honor j


class OpaqueNodeError(AdapterError):
    """Raised when an op targets an out-of-scope (opaque, untracked) shape."""


@dataclass
class _Paragraph:
    """One tracked paragraph inside a text frame (text + optional style hint)."""

    text: str
    style: Optional[str] = None


@dataclass
class _Shape:
    """Adapter-side state for one shape on a slide (opaque shapes carry a label)."""

    shape_id: str
    name: str
    opaque: bool
    paragraphs: list[_Paragraph] = field(default_factory=list)
    label: str = ""


@dataclass
class _Slide:
    """Adapter-side state for one slide, keyed by its native ``slide_id``."""

    slide_id: str
    layout: str
    title: str
    shapes: list[_Shape] = field(default_factory=list)
    inserted: bool = False  # added via slide.insert (no live pptx slide yet)


def _has_text_frame(shape: Any) -> bool:
    try:
        return bool(shape.has_text_frame)
    except (AttributeError, ValueError):  # pragma: no cover - defensive
        return False


def _shape_label(shape: Any) -> str:
    """Return a coarse kind label for an opaque (out-of-scope) shape."""
    if getattr(shape, "has_chart", False):
        return "chart"
    if getattr(shape, "has_table", False):
        return "table"
    lowered = str(getattr(shape, "name", "") or "").lower()
    for kind in ("picture", "image", "smartart", "diagram", "group", "media"):
        if kind in lowered:
            return kind
    return lowered or "shape"


def _in_scope(shape: Any) -> bool:
    return (
        _has_text_frame(shape)
        and not getattr(shape, "has_chart", False)
        and not getattr(shape, "has_table", False)
    )


class PptxAdapter(DocumentAdapter):
    """Loads a .pptx, applies the in-scope v0.2 slide/text ops, renders bytes."""

    def __init__(
        self,
        raw_bytes: bytes,
        *,
        author: str = DEFAULT_AUTHOR,
        date: str = DEFAULT_DATE,
    ) -> None:
        self._raw = raw_bytes
        self._author = author
        self._date = date
        self._baseline_sha = sha256_hex(raw_bytes)
        self._minted = 0
        self._slides: list[_Slide] = []
        self._build_model_from_bytes(raw_bytes)

    @classmethod
    def load(
        cls,
        path: str,
        *,
        author: str = DEFAULT_AUTHOR,
        date: str = DEFAULT_DATE,
        **_ignored: object,
    ) -> "PptxAdapter":
        """Load a .pptx from a sanitized path (extra kwargs accepted+ignored)."""
        resolved = safe_path(path, must_exist=True, allow_suffixes=(".pptx",))
        return cls(resolved.read_bytes(), author=author, date=date)

    def _build_model_from_bytes(self, raw: bytes) -> None:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(raw))
        self._slides = [self._slide_state(s) for s in prs.slides]

    def _slide_state(self, slide: Any) -> _Slide:
        layout = str(getattr(getattr(slide, "slide_layout", None), "name", "") or "")
        title = ""
        try:
            if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
                title = str(slide.shapes.title.text_frame.text)
        except (AttributeError, ValueError):
            title = ""
        shapes = [self._shape_state(sh) for sh in slide.shapes]
        return _Slide(
            slide_id=str(slide.slide_id), layout=layout, title=title, shapes=shapes
        )

    def _shape_state(self, shape: Any) -> _Shape:
        shape_id = str(getattr(shape, "shape_id", ""))
        name = str(getattr(shape, "name", "") or "")
        if not _in_scope(shape):
            return _Shape(shape_id, name, opaque=True, label=_shape_label(shape))
        paras = [_Paragraph(str(p.text)) for p in shape.text_frame.paragraphs]
        return _Shape(shape_id, name, opaque=False, paragraphs=paras)

    def _mint_slide_id(self) -> str:
        self._minted += 1  # synthetic id for a not-yet-live slide
        return f"new{self._minted:04d}"

    def baseline_sha256(self) -> str:
        return self._baseline_sha

    def to_model(self) -> Node:
        """Return the current tree (slides -> shapes -> paragraphs), positionally."""
        root = Node(node_id="root", node_kind=NodeKind.DOCUMENT, path="/presentation")
        for i, slide in enumerate(self._slides):
            slide_node = Node(
                node_id=f"slide[{i}]",
                node_kind=NodeKind.DOCUMENT,
                path=f"/presentation/slide[{i + 1}]",
                value=slide.title,
                attrs={
                    "kind": "slide", "slide_id": slide.slide_id,
                    "layout": slide.layout, "title": slide.title,
                    "inserted": slide.inserted,
                },
            )
            for shape in slide.shapes:
                slide_node.children.append(self._shape_node(i, shape, slide_node.path))
            root.children.append(slide_node)
        return root

    def _shape_node(self, slide_idx: int, shape: _Shape, slide_path: str) -> Node:
        addr = f"slide[{slide_idx}]/shape:{shape.shape_id}"
        path = f"{slide_path}/shape:{shape.shape_id}"
        attrs = {"kind": "shape", "opaque": shape.opaque, "name": shape.name}
        if shape.opaque:
            attrs["shape_kind"] = shape.label
        node = Node(node_id=addr, node_kind=NodeKind.PARAGRAPH, path=path, attrs=attrs)
        for j, para in enumerate(shape.paragraphs):
            node.children.append(
                Node(
                    node_id=f"{addr}/p[{j}]",
                    node_kind=NodeKind.RUN,
                    path=f"{path}/p[{j + 1}]",
                    value=para.text,
                    attrs={
                        "kind": "paragraph",
                        **({"style": para.style} if para.style else {}),
                    },
                )
            )
        return node

    def set_model(self, root: Node) -> None:
        """Reset state to ``root`` (used by replay), keyed on recorded slide_ids."""
        self._slides = []
        for slide_node in root.children:
            a = slide_node.attrs
            slide = _Slide(
                slide_id=str(a.get("slide_id", self._mint_slide_id())),
                layout=str(a.get("layout", "")),
                title=str(a.get("title", slide_node.value or "")),
                inserted=bool(a.get("inserted", False)),
            )
            for shape_node in slide_node.children:
                slide.shapes.append(self._shape_from_node(shape_node))
            self._slides.append(slide)

    @staticmethod
    def _shape_from_node(shape_node: Node) -> _Shape:
        a = shape_node.attrs
        shape_id = shape_node.node_id.rsplit("shape:", 1)[-1].split("/", 1)[0]
        if bool(a.get("opaque", False)):
            return _Shape(
                shape_id,
                str(a.get("name", "")),
                opaque=True,
                label=str(a.get("shape_kind", "shape")),
            )
        paras = [
            _Paragraph(
                str(p.value or ""),
                style=str(p.attrs["style"]) if p.attrs.get("style") else None,
            )
            for p in shape_node.children
        ]
        return _Shape(shape_id, str(a.get("name", "")), opaque=False, paragraphs=paras)

    def resolve(self, node_id: str) -> Node | None:
        return self.to_model().find(node_id)

    def apply(self, op: Op) -> None:
        """Apply one in-scope v0.2 op (slide.insert/delete or shape.edit).

        Raises BeforeMismatchError/OversizedOpError on bad text ops,
        NodeNotFoundError on a missing slide/shape/paragraph, OpaqueNodeError on
        an out-of-scope shape, and TypeError for any other op kind.
        """
        if isinstance(op, SlideInsert):
            self._apply_slide_insert(op)
        elif isinstance(op, SlideDelete):
            self._apply_slide_delete(op)
        elif isinstance(op, ShapeEdit):
            self._apply_shape_edit(op)
        else:  # pragma: no cover - exhaustive over the pptx scope
            raise TypeError(
                f"pptx adapter does not handle op {type(op).__name__}; supported: "
                "slide.insert, slide.delete, shape.edit"
            )

    def _apply_slide_insert(self, op: SlideInsert) -> None:
        title = str(op.value.get("title", ""))
        new_slide = _Slide(
            slide_id=self._mint_slide_id(),
            layout=str(op.value.get("layout", "")),
            title=title,
            inserted=True,
        )
        if title:
            new_slide.shapes.append(
                _Shape("title", "Title", opaque=False, paragraphs=[_Paragraph(title)])
            )
        pos = max(0, min(int(op.at), len(self._slides)))
        self._slides.insert(pos, new_slide)

    def _apply_slide_delete(self, op: SlideDelete) -> None:
        at = int(op.at)
        if at < 0 or at >= len(self._slides):
            raise NodeNotFoundError(f"no slide at position {at}")
        del self._slides[at]

    def _apply_shape_edit(self, op: ShapeEdit) -> None:
        idx = int(op.slide)
        if idx < 0 or idx >= len(self._slides):
            raise NodeNotFoundError(f"no slide at position {idx}")
        shape = self._find_shape(self._slides[idx], op.shape_id)
        if shape.opaque:
            raise OpaqueNodeError(
                f"shape {op.shape_id!r} on slide {idx} is out of scope "
                f"({shape.label}); charts/tables/pictures/SmartArt are not tracked"
            )
        self._apply_text_op(op.op, self._resolve_para(op.op, shape))

    @staticmethod
    def _find_shape(slide: _Slide, shape_id: str) -> _Shape:
        for shape in slide.shapes:
            if shape.shape_id == shape_id:
                return shape
        raise NodeNotFoundError(
            f"no shape with id {shape_id!r} on slide {slide.slide_id}"
        )

    @staticmethod
    def _resolve_para(raw_op: dict[str, Any], shape: _Shape) -> _Paragraph:
        """Locate the target paragraph; honors only the node_id's ``p[j]`` index."""
        match = _PARA_NODE_RE.match(str(raw_op.get("node_id", "") or ""))
        idx = int(match.group(3)) if match else 0
        if not shape.paragraphs:
            shape.paragraphs.append(_Paragraph(""))
        if idx < 0 or idx >= len(shape.paragraphs):
            raise NodeNotFoundError(f"no paragraph p[{idx}] in shape {shape.shape_id!r}")
        return shape.paragraphs[idx]

    def _check_size(self, para: _Paragraph, changed_len: int) -> None:
        current = len(para.text)
        if current and changed_len > current * MAX_OP_FRACTION:
            pct = int(MAX_OP_FRACTION * 100)
            raise OversizedOpError(
                f"split_required: this op rewrites more than {pct}% of the "
                "paragraph; split it into smaller edits and re-issue them."
            )

    def _apply_text_op(self, raw_op: dict[str, Any], para: _Paragraph) -> None:
        """Apply a nested text/style sub-op (the only kinds shape.edit supports)."""
        kind = str(raw_op.get("kind", ""))
        if kind == "text.replace":
            self._text_replace(para, str(raw_op["before"]), str(raw_op["after"]))
        elif kind == "text.insert":
            self._text_insert(para, raw_op.get("before_anchor"), str(raw_op["text"]))
        elif kind == "text.delete":
            self._text_delete(para, str(raw_op["before"]))
        elif kind == "style.change":
            self._style_change(
                para, str(raw_op["style"]), str(raw_op.get("before", ""))
            )
        else:
            raise TypeError(
                f"shape.edit nested op {kind!r} is out of scope; supported: "
                "text.insert, text.delete, text.replace, style.change"
            )

    def _text_replace(self, para: _Paragraph, before: str, after: str) -> None:
        self._replace_first(para, before, after, max(len(before), len(after)))

    def _text_delete(self, para: _Paragraph, before: str) -> None:
        self._replace_first(para, before, "", len(before))

    def _replace_first(
        self, para: _Paragraph, before: str, after: str, changed_len: int
    ) -> None:
        """Validate ``before`` (non-empty + present) and replace its first hit."""
        if before == "":
            raise BeforeMismatchError("before must be non-empty")
        self._check_size(para, changed_len)
        if before not in para.text:
            raise BeforeMismatchError(f"before text {before!r} not in paragraph")
        para.text = para.text.replace(before, after, 1)

    def _text_insert(
        self, para: _Paragraph, before_anchor: Optional[str], text: str
    ) -> None:
        self._check_size(para, len(text))
        if not before_anchor:
            para.text = para.text + text
            return
        if before_anchor not in para.text:
            raise BeforeMismatchError(f"anchor {before_anchor!r} not in paragraph")
        idx = para.text.find(before_anchor) + len(before_anchor)
        para.text = para.text[:idx] + text + para.text[idx:]

    @staticmethod
    def _style_change(para: _Paragraph, style: str, before: str) -> None:
        current = para.style or ""
        if before and before != current:
            raise BeforeMismatchError(f"style before {before!r} != current {current!r}")
        para.style = style

    def render_tracked(self) -> bytes:
        """Return the **edited deck** bytes (pptx has no native track-changes).

        Replays the model's slide/text state onto a fresh copy of the baseline (no
        in-file revision markup — the journal is authoritative; see
        :func:`changex_core.render.pptx.render_review_overlay`).
        """
        from pptx import Presentation

        prs = Presentation(io.BytesIO(self._raw))
        # model slide_id -> live pptx slide (existing by native id; inserted minted).
        model_to_live: dict[str, Any] = {}
        live_by_native = {str(s.slide_id): s for s in prs.slides}
        for slide_state in self._slides:
            live = live_by_native.get(slide_state.slide_id)
            if live is None and slide_state.inserted:
                live = self._append_pptx_slide(prs, slide_state)
            if live is not None:
                self._render_text(live, slide_state)
                model_to_live[slide_state.slide_id] = live
        self._reorder_slides(prs, model_to_live)
        buffer = io.BytesIO()
        prs.save(buffer)
        return buffer.getvalue()

    def _render_text(self, live_slide: Any, slide_state: _Slide) -> None:
        live_shapes = {str(sh.shape_id): sh for sh in live_slide.shapes}
        for shape_state in slide_state.shapes:
            if shape_state.opaque:
                continue
            live = live_shapes.get(shape_state.shape_id)
            if live is not None and _has_text_frame(live):
                self._write_text_frame(live.text_frame, shape_state.paragraphs)

    def _write_text_frame(self, text_frame: Any, paragraphs: list[_Paragraph]) -> None:
        existing = list(text_frame.paragraphs)
        for i, para in enumerate(paragraphs):
            if i < len(existing):
                self._set_paragraph_text(existing[i], para.text)
            else:
                self._set_paragraph_text(text_frame.add_paragraph(), para.text)
        for extra in existing[len(paragraphs):]:
            self._set_paragraph_text(extra, "")

    @staticmethod
    def _set_paragraph_text(paragraph: Any, text: str) -> None:
        """Set a paragraph's text via its first run (preserving formatting)."""
        runs = list(paragraph.runs)
        if runs:
            runs[0].text = text
            for extra in runs[1:]:
                extra.text = ""
        elif text:
            paragraph.add_run().text = text

    def _append_pptx_slide(self, prs: Any, slide_state: _Slide) -> Any:
        slide = prs.slides.add_slide(self._pick_layout(prs, slide_state.layout))
        if slide_state.title:
            if slide.shapes.title is not None:
                slide.shapes.title.text = slide_state.title
            else:  # layout lacks a title placeholder -> carry the title in a textbox
                from pptx.util import Inches

                box = slide.shapes.add_textbox(*(Inches(x) for x in (0.5, 0.3, 9, 1)))
                box.text_frame.text = slide_state.title
        return slide

    @staticmethod
    def _pick_layout(prs: Any, layout_name: str) -> Any:
        for layout in prs.slide_layouts:
            if str(layout.name) == layout_name:
                return layout
        return prs.slide_layouts[min(_BLANK_LAYOUT_INDEX, len(prs.slide_layouts) - 1)]

    def _reorder_slides(self, prs: Any, model_to_live: dict[str, Any]) -> None:
        """Reorder/remove pptx slides to match the model's slide sequence.

        Pair each live slide with its ``sldId`` element BEFORE any detach (a
        detached slide cannot resolve its id).
        """
        xml_slides = prs.slides._sldIdLst
        sld_elements = list(xml_slides)
        # live-slide identity -> its sldId element, then model slide_id -> element.
        el_by_live = {id(s): el for s, el in zip(list(prs.slides), sld_elements)}
        model_to_el = {
            mid: el_by_live[id(live)]
            for mid, live in model_to_live.items()
            if id(live) in el_by_live
        }
        for el in sld_elements:
            xml_slides.remove(el)
        for slide in self._slides:
            el = model_to_el.get(slide.slide_id)
            if el is not None:
                xml_slides.append(el)

    def save(self, out_path: str) -> None:
        """Save the **edited deck** bytes (not a tracked-changes file: pptx has
        none; see :func:`changex_core.render.pptx.save_review_overlay`)."""
        resolved = safe_path(out_path, allow_suffixes=(".pptx",))
        resolved.parent.mkdir(parents=True, exist_ok=True)
        resolved.write_bytes(self.render_tracked())

    def node_id_map(self) -> dict[str, str]:
        """Return a ``node_id -> carrier`` map (positional slide -> native id)."""
        return {f"slide[{i}]": s.slide_id for i, s in enumerate(self._slides)}
